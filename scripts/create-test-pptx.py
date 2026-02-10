"""
创建简单的 PPTX 测试文件
运行: python scripts/create-test-pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import sys
import io

# 设置输出编码为 UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

def create_test_pptx():
    """创建一个简单的测试用 PPTX 文件"""

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 第1张幻灯片 - 标题页
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "PPTX to JSON Conversion Test"
    subtitle.text = "End-to-End Integration Test"

    # 第2张幻灯片 - 包含多种元素
    content_slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 添加标题
    title_shape = content_slide.shapes.title
    title_shape.text = "Test Elements"

    # 添加文本框
    left = Inches(1)
    top = Inches(2)
    width = Inches(3)
    height = Inches(1)
    textbox = content_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.text = "This is a test text box"
    p.font.size = Pt(18)

    # 添加矩形
    left = Inches(5)
    top = Inches(2)
    width = Inches(2)
    height = Inches(1.5)
    shape = content_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left,
        top,
        width,
        height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x70, 0xC0, 0x00)  # 蓝色
    shape.line.color.rgb = RGBColor(0, 0, 0)  # 黑色边框
    shape.line.width = Pt(1)

    # 添加另一个文本框
    left = Inches(1)
    top = Inches(4)
    width = Inches(6)
    height = Inches(1)
    textbox2 = content_slide.shapes.add_textbox(left, top, width, height)
    text_frame2 = textbox2.text_frame
    p2 = text_frame2.paragraphs[0]
    p2.text = "Element 2: Another text box"
    p2.font.size = Pt(14)
    p2.font.bold = True

    # 保存文件
    output_path = 'tests/fixtures/simple.pptx'
    prs.save(output_path)

    print(f"[OK] Successfully created test PPTX file: {output_path}")
    print(f"   Slide count: {len(prs.slides)}")
    print(f"   File size: ~{len(prs.slides)} slides with basic elements")

    return output_path

if __name__ == "__main__":
    try:
        create_test_pptx()
    except ImportError:
        print("[ERROR] Need to install python-pptx")
        print("   Run: pip install python-pptx")
    except Exception as e:
        print(f"[ERROR] Failed to create file: {e}")
        import traceback
        traceback.print_exc()

