#!/usr/bin/env python3
"""
创建一个简单的 PPTX 测试文件
"""
from pptx import Presentation
from pptx.util import Inches, Pt

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 添加标题幻灯片
slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
slide = prs.slides.add_slide(slide_layout)

# 设置标题
title = slide.shapes.title
title.text = "测试演示文稿"

# 设置副标题
subtitle = slide.placeholders[1]
subtitle.text = "这是一个简单的测试 PPTX"

# 添加内容幻灯片
bullet_slide_layout = prs.slide_layouts[1]  # 项目符号布局
slide = prs.slides.add_slide(bullet_slide_layout)

# 添加标题
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = "主要内容"

# 添加内容文本框
left = Inches(1)
top = Inches(2)
width = Inches(8)
height = Inches(4)

textbox = shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# 添加段落
p = text_frame.paragraphs[0]
p.text = "这是第一行文本"
p.font.size = Pt(18)

p = text_frame.add_paragraph()
p.text = "这是第二行文本"
p.font.size = Pt(18)
p.level = 1

p = text_frame.add_paragraph()
p.text = "这是第三行文本"
p.font.size = Pt(18)

# 保存文件
prs.save('simple.pptx')
print("simple.pptx created successfully")
